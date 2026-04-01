from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BRAND   = RGBColor(0x63, 0x66, 0xF1)
DARK    = RGBColor(0x0F, 0x10, 0x1A)
SURFACE = RGBColor(0x1A, 0x1C, 0x2E)
MUTED   = RGBColor(0x6B, 0x72, 0x80)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER   = RGBColor(0xF5, 0x9E, 0x0B)
VIOLET  = RGBColor(0x8B, 0x5C, 0xF6)
ROSE    = RGBColor(0xF4, 0x3F, 0x5E)
CYAN    = RGBColor(0x06, 0xB6, 0xD4)
PINK    = RGBColor(0xEC, 0x48, 0x99)
TEAL    = RGBColor(0x14, 0xB8, 0xA6)
INDIGO  = RGBColor(0xA5, 0xB4, 0xFC)

blank_layout = prs.slide_layouts[6]

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill_color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def accent_bar(slide, color=BRAND):
    rect(slide, 0, 0, 0.22, 7.5, color)

def section_header(slide, number, label):
    rect(slide, 0, 0, 13.33, 1.1, SURFACE)
    txt(slide, f"{number:02d}", 0.4, 0.08, 1.0, 0.9, size=44, bold=True, color=BRAND)
    txt(slide, label, 1.55, 0.22, 11.0, 0.7, size=26, bold=True, color=WHITE)

def card(slide, l, t, w, h, accent=BRAND):
    rect(slide, l, t, w, h, SURFACE)
    rect(slide, l, t, 0.06, h, accent)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 01 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, 0, 0, 4.6, 7.5, SURFACE)
rect(s, 0, 0, 0.08, 7.5, BRAND)

txt(s, "KAIROS", 0.5, 1.6, 4.0, 1.0, size=52, bold=True, color=WHITE)
txt(s, "Consulting Management Platform", 0.5, 2.7, 4.0, 0.45, size=14, color=MUTED)
rect(s, 0.5, 3.25, 3.5, 0.055, BRAND)
txt(s, "Track  ·  Review  ·  Invoice  ·  Grow", 0.5, 3.42, 4.0, 0.4,
    size=12, italic=True, color=INDIGO)
txt(s, "From first logged hour to compliant invoice —\nKairos manages the full consulting workflow\nin one role-aware platform.",
    0.5, 4.05, 4.0, 1.1, size=11, color=MUTED)
txt(s, "Product Overview  ·  2025", 0.5, 6.9, 4.0, 0.38, size=10, color=MUTED)

# right hero stats
stats = [
    ("9",    "Modules",     BRAND,   5.0, 1.3),
    ("4",    "User Roles",  VIOLET,  9.0, 1.3),
    ("EN\n16931", "Compliant", EMERALD, 5.0, 3.5),
    ("DE/EN","Bilingual",   AMBER,   9.0, 3.5),
]
for val, label, col, x, y in stats:
    card(s, x, y, 3.5, 1.85, col)
    txt(s, val,   x+0.2, y+0.2,  3.0, 0.9,  size=32, bold=True, color=col)
    txt(s, label, x+0.2, y+1.2,  3.0, 0.45, size=12, color=MUTED)

txt(s, '"The only tool your consulting firm needs — from time entry to BMD export."',
    4.8, 6.1, 8.3, 0.8, size=11, italic=True, color=INDIGO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 02 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "The Problem", 0.55, 0.28, 9, 0.6, size=30, bold=True, color=WHITE)
txt(s, "Consulting firms lose revenue every day because of fragmented tooling",
    0.55, 0.93, 12.5, 0.38, size=13, color=MUTED)

pains = [
    ("Hours leak away",    "Consultants forget to track, or log from memory days later. Billable time silently disappears.", ROSE),
    ("No approval chain",  "Hours are invoiced before a PM or partner has reviewed them — disputes and write-offs follow.", AMBER),
    ("Invoice chaos",      "Finance manually assembles invoices from spreadsheets. Errors, delays, and compliance gaps.", ROSE),
    ("Zero visibility",    "Management has no real-time view of utilization or budget status until the month is over.", AMBER),
    ("Tool sprawl",        "Teams juggle Toggl + Jira + Excel + Outlook + Datev — no single source of truth.", ROSE),
    ("Compliance gaps",    "Austrian firms need EN 16931 invoices & BMD NTCS exports — unavailable in generic tools.", AMBER),
]
for i, (title, body, col) in enumerate(pains):
    c = i % 2
    r = i // 2
    x = 0.35 + c * 6.55
    y = 1.55 + r * 1.85
    card(s, x, y, 6.1, 1.68, col)
    txt(s, title, x+0.2, y+0.12, 5.7, 0.4,  size=13, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.6,  5.7, 0.95, size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 03 — SOLUTION / MODULE MAP
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "One Platform. Every Workflow.", 0.55, 0.28, 12, 0.6, size=30, bold=True, color=WHITE)
txt(s, "9 integrated modules — no stitching required",
    0.55, 0.93, 12, 0.38, size=13, color=MUTED)

modules = [
    ("Timer",       "Log hours: live, manual, or duration",   BRAND),
    ("Dashboard",   "Personal & team KPIs at a glance",        VIOLET),
    ("Timesheets",  "Submit, review, approve per project",      EMERALD),
    ("Projects",    "Budgets, rates, team, PM assignment",      AMBER),
    ("Clients",     "Client CRM with logo & billing info",      ROSE),
    ("Invoices",    "Generate & export EN 16931 invoices",      CYAN),
    ("Analytics",   "6-month revenue & utilization trends",     PINK),
    ("Reports",     "Custom date-range reports + CSV export",   TEAL),
    ("Settings",    "Team, levels, workspace, compliance",      MUTED),
]
for i, (label, sub, col) in enumerate(modules):
    c = i % 3
    r = i // 3
    x = 0.35 + c * 4.32
    y = 1.55 + r * 1.85
    card(s, x, y, 4.0, 1.68, col)
    txt(s, label, x+0.2, y+0.18, 3.6, 0.5,  size=16, bold=True, color=WHITE)
    txt(s, sub,   x+0.2, y+0.78, 3.6, 0.72, size=11, color=MUTED)

# flow bar
rect(s, 0.35, 7.08, 12.6, 0.32, RGBColor(0x12, 0x13, 0x24))
flow = "Log Hours  →  Submit Timesheet  →  PM / Partner Review  →  Generate Invoice  →  BMD Export"
txt(s, flow, 0.6, 7.1, 12.2, 0.28, size=11, color=INDIGO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 04 — ROLES
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Built for Every Role in the Firm", 0.55, 0.28, 12, 0.6, size=30, bold=True, color=WHITE)
txt(s, "Four distinct roles — each with precisely scoped permissions",
    0.55, 0.93, 12.5, 0.38, size=13, color=MUTED)

roles = [
    ("Admin", ROSE, [
        "Full workspace control",
        "Invite & manage team",
        "Approve all timesheets",
        "Create & send invoices",
        "View full analytics",
        "Does NOT track time",
    ]),
    ("Partner", AMBER, [
        "Tracks billable time",
        "Reviews all timesheets",
        "Manages clients & projects",
        "Creates invoices",
        "Full analytics access",
        "",
    ]),
    ("Project Manager", BRAND, [
        "Tracks billable time",
        "Approves own project hours",
        "Manages assigned projects",
        "Sets per-level rates",
        "Views project analytics",
        "",
    ]),
    ("Consultant / Member", EMERALD, [
        "Logs time entries",
        "Submits weekly timesheets",
        "Views own hours & reports",
        "Cannot see team data",
        "Cannot manage projects",
        "",
    ]),
]
for i, (role, col, perms) in enumerate(roles):
    x = 0.35 + i * 3.22
    rect(s, x, 1.55, 3.0, 0.55, col)
    txt(s, role, x, 1.58, 3.0, 0.5, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    rect(s, x, 2.1, 3.0, 5.1, SURFACE)
    for j, p in enumerate(perms):
        if p:
            txt(s, f"✓  {p}", x+0.15, 2.25+j*0.78, 2.75, 0.65, size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 05 — TIME TRACKING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 1, "Time Tracking  ·  Three Entry Modes")
txt(s, "Live timer · manual from/to · duration entry · automatic deadline enforcement",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

modes = [
    ("Live Timer",       "Start/stop with one click. Real-time elapsed display. Idle alert after 3 hours continuous tracking.", BRAND),
    ("From  →  To",      "Enter start and end time. Ideal for logging meetings or past work sessions precisely.", VIOLET),
    ("Duration Entry",   "Type hours directly. Fastest method for fixed-length tasks like 4h workshops.", EMERALD),
]
for i, (title, body, col) in enumerate(modes):
    card(s, 0.35, 1.68+i*1.68, 5.6, 1.52, col)
    txt(s, title, 0.6, 1.83+i*1.68, 5.0, 0.42, size=13, bold=True, color=WHITE)
    txt(s, body,  0.6, 2.3+i*1.68,  5.2, 0.72, size=11, color=MUTED)

callouts = [
    ("Idle Detection",    "Alert after 3 h — trim, keep, or discard the entry.",       BRAND),
    ("Auto-Lock",         "Entries lock after Sunday 23:00 deadline. Admins unlock.",   ROSE),
    ("Billable Flag",     "Per-entry billable toggle. Rate snapshot at entry time.",     EMERALD),
    ("Time Rounding",     "5 / 10 / 15 / 30 min or 1 h rounding — per project.",       AMBER),
    ("Team Monitor",      "Admins see who is tracking now, how long, on what project.", VIOLET),
    ("Full Edit Trail",   "All edits tracked. History visible even after submission.",  MUTED),
]
for i, (title, body, col) in enumerate(callouts):
    c = i % 2
    r = i // 2
    x = 6.3 + c * 3.45
    y = 1.68 + r * 1.68
    card(s, x, y, 3.2, 1.52, col)
    txt(s, title, x+0.2, y+0.12, 2.8, 0.42, size=12, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.62, 2.8, 0.75, size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 06 — TIMESHEET APPROVAL
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 2, "Timesheet Approval  ·  Per-Project Authority")
txt(s, "The right person approves the right hours — enforced in UI and server-side",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

statuses = [
    ("Draft",     MUTED,    "Week in progress"),
    ("Submitted", AMBER,    "Sunday deadline"),
    ("Reviewing", BRAND,    "PM or Partner"),
    ("Approved",  EMERALD,  "Ready to invoice"),
    ("Rejected",  ROSE,     "Returned + note"),
]
for i, (label, col, sub) in enumerate(statuses):
    x = 0.35 + i * 2.57
    rect(s, x, 1.55, 2.3, 0.62, col)
    txt(s, label, x, 1.57, 2.3, 0.55, size=13, bold=True,
        color=DARK if col==AMBER else WHITE, align=PP_ALIGN.CENTER)
    txt(s, sub, x+0.05, 2.22, 2.2, 0.35, size=9, color=MUTED, align=PP_ALIGN.CENTER)
    if i < 4:
        txt(s, "→", x+2.34, 1.68, 0.3, 0.38, size=14, bold=True, color=BRAND, align=PP_ALIGN.CENTER)

# Per-project approval block
rect(s, 0.35, 2.75, 12.6, 4.4, SURFACE)
txt(s, "Per-Project Approval — How It Works", 0.6, 2.9, 9, 0.42, size=14, bold=True, color=WHITE)
txt(s, "A PM can only approve hours belonging to their own projects. "
       "Partners and Admins approve the full timesheet.",
    0.6, 3.38, 12.0, 0.35, size=11, color=MUTED)

steps = [
    ("Member A submits week",       "5 h on Project Alpha  +  2 h on Project Beta",         WHITE,   SURFACE),
    ("PM of Project Alpha reviews", "Sees Project Alpha hours only → Approves → Alpha ✓",    BRAND,   SURFACE),
    ("PM of Project Beta reviews",  "Sees Project Beta hours only  → Approves → Beta ✓",     BRAND,   SURFACE),
    ("System auto-approves",        "All projects signed off → timesheet status = Approved",  EMERALD, SURFACE),
    ("Admin override available",    "Admin or Partner can approve the entire timesheet at any time", AMBER, SURFACE),
]
for i, (actor, action, acol, _) in enumerate(steps):
    y = 3.88 + i * 0.58
    rect(s, 0.5, y, 0.42, 0.42, acol)
    txt(s, str(i+1), 0.5, y, 0.42, 0.4, size=12, bold=True,
        color=DARK, align=PP_ALIGN.CENTER)
    txt(s, actor,  1.1, y+0.04, 5.0, 0.35, size=11, bold=True, color=acol)
    txt(s, action, 6.3, y+0.04, 6.4, 0.35, size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 07 — PROJECT MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 3, "Project Management  ·  Budgets · Rates · Team")
txt(s, "Consultant level rate matrix · real-time budget warnings · PM-scoped approval authority",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

features = [
    ("Level Rate Matrix",   "Set different hourly or daily rates per consultant level per project. Junior ≠ Senior ≠ Partner.", BRAND),
    ("Budget Tracking",     "Define hour & amount budgets. Warning at 80 %, red at 100 %. Analytics predicts exhaustion date.", AMBER),
    ("Team Assignment",     "Assign consultants to projects. Members only see assigned projects — others stay hidden.", EMERALD),
    ("PM Assignment",       "Designate one PM per project. PM gets timesheet review authority scoped to that project only.", VIOLET),
    ("Time Rounding",       "Configure 5 / 10 / 15 / 30 min or 1 h auto-rounding rules per project.", ROSE),
    ("Archive & Preserve",  "Archive completed projects. Historical data intact. Soft-delete keeps records accessible.", MUTED),
]
for i, (title, body, col) in enumerate(features):
    c = i % 2
    r = i // 2
    x = 0.35 + c * 6.55
    y = 1.55 + r * 1.9
    card(s, x, y, 6.1, 1.72, col)
    txt(s, title, x+0.2, y+0.12, 5.7, 0.42, size=13, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.65, 5.7, 0.9,  size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 08 — INVOICING
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 4, "Invoicing  ·  EN 16931 · BMD NTCS · Compliant")
txt(s, "From approved hours to tax-compliant invoice in minutes — with direct accounting export",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

# Left: steps
rect(s, 0.35, 1.6, 5.7, 5.7, SURFACE)
txt(s, "Invoice Generation Flow", 0.6, 1.75, 5.2, 0.4, size=13, bold=True, color=WHITE)
steps = [
    ("1", "Select client & billing date range",     BRAND),
    ("2", "System pulls approved hours only",        EMERALD),
    ("3", "Choose VAT code, set due date & notes",   AMBER),
    ("4", "Assign invoice number, preview",          BRAND),
    ("5", "Export PDF or BMD NTCS file",             VIOLET),
    ("6", "Mark as Sent → track until Paid",         EMERALD),
]
for i, (num, step, col) in enumerate(steps):
    y = 2.28 + i * 0.8
    rect(s, 0.6, y, 0.45, 0.45, col)
    txt(s, num,  0.6, y, 0.45, 0.42, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, step, 1.2, y+0.04, 4.6, 0.38, size=12, color=WHITE)

# Right top: VAT
rect(s, 6.4, 1.6, 6.6, 2.65, SURFACE)
txt(s, "Austrian VAT Codes Supported", 6.65, 1.75, 6.1, 0.38, size=12, bold=True, color=WHITE)
vat = [
    ("U20", "20 % Standard (Inland)",           BRAND),
    ("U10", "10 % Reduced rate",                VIOLET),
    ("IG",  "0 % EU Services (§3a UStG)",       EMERALD),
    ("RC",  "0 % Reverse Charge (§19 UStG)",    AMBER),
    ("AU",  "0 % Export / Ausfuhrlieferung",    MUTED),
]
for i, (code, label, col) in enumerate(vat):
    y = 2.22 + i * 0.4
    rect(s, 6.65, y, 0.62, 0.32, col)
    txt(s, code, 6.65, y, 0.62, 0.3, size=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, label, 7.42, y+0.02, 5.3, 0.28, size=10, color=WHITE)

# Right bottom: compliance
rect(s, 6.4, 4.4, 6.6, 2.9, SURFACE)
txt(s, "Compliance & Audit", 6.65, 4.55, 6.1, 0.38, size=12, bold=True, color=WHITE)
comp = [
    "Only approved timesheet hours invoiceable",
    "Seller & buyer details snapshotted at issue",
    "EN 16931-compliant invoice structure",
    "BMD NTCS export (configurable accounts)",
    "Erlöskonto + Debitorenkonto mapping",
    "Draft → Sent → Paid status lifecycle",
]
for i, line in enumerate(comp):
    txt(s, f"✓  {line}", 6.65, 5.05+i*0.36, 6.1, 0.32, size=10, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 09 — ANALYTICS & DASHBOARD
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 5, "Analytics & Dashboard  ·  Real-Time Intelligence")
txt(s, "Personal KPIs for every role · 6-month team analytics · burnout risk detection",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

txt(s, "Personal Dashboard  (all roles)", 0.45, 1.6, 7, 0.35, size=11, bold=True, color=MUTED)
dash = [
    ("This Week",    "Hours + trend vs prev. week",  BRAND),
    ("This Month",   "Cumulative hours MTD",          VIOLET),
    ("Earnings",     "Billable revenue this month",   EMERALD),
    ("Utilization",  "Billable % of total capacity",  AMBER),
]
for i, (title, sub, col) in enumerate(dash):
    x = 0.35 + i * 3.22
    card(s, x, 2.0, 3.0, 1.4, col)
    txt(s, title, x+0.2, 2.15, 2.6, 0.42, size=13, bold=True, color=WHITE)
    txt(s, sub,   x+0.2, 2.62, 2.6, 0.62, size=10, color=MUTED)

txt(s, "Admin / Partner Analytics  (6-month view)", 0.45, 3.55, 9, 0.35, size=11, bold=True, color=MUTED)
analytics = [
    ("Revenue MTD",       "Month-to-date with trend vs. prior month",                    BRAND),
    ("Budget Burndown",   "Remaining budget + projected exhaustion date",                 ROSE),
    ("Team Utilization",  "Per-member billable % · burnout flag at >100 % for 3+ weeks", AMBER),
    ("Revenue by Client", "Recharts pie chart — which clients drive the most revenue",    VIOLET),
    ("Invoice Pipeline",  "Collected / Outstanding / Overdue amounts",                   EMERALD),
    ("Week-by-Week Trend","24-week utilization area chart for capacity planning",         CYAN),
]
for i, (title, body, col) in enumerate(analytics):
    c = i % 3
    r = i // 3
    x = 0.35 + c * 4.32
    y = 4.05 + r * 1.55
    card(s, x, y, 4.0, 1.38, col)
    txt(s, title, x+0.2, y+0.1,  3.6, 0.42, size=12, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.6,  3.6, 0.65, size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — REPORTS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 6, "Reports  ·  Custom Ranges · Three Views · CSV Export")
txt(s, "Drill into hours, earnings, and team performance — export raw data in one click",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

views = [
    ("Overview",     "Total hours · non-billable hours · billable ratio %\nTotal revenue · daily bar chart",        BRAND),
    ("Team View",    "Per-member breakdown: hours, billable %, earnings\nProject time distribution per consultant", VIOLET),
    ("Entry Detail", "Every entry with full metadata:\ndate, project, client, description, billable flag, rate",    EMERALD),
]
for i, (title, body, col) in enumerate(views):
    x = 0.35 + i * 4.32
    card(s, x, 1.6, 4.0, 2.35, col)
    txt(s, title, x+0.2, 1.75, 3.6, 0.45, size=15, bold=True, color=WHITE)
    txt(s, body,  x+0.2, 2.32, 3.6, 1.45, size=11, color=MUTED)

rect(s, 0.35, 4.1, 12.6, 3.05, SURFACE)
txt(s, "Filters & Capabilities", 0.6, 4.25, 6, 0.4, size=13, bold=True, color=WHITE)
feats = [
    ("Date Range Presets",  "This/last week  ·  this/last month  ·  custom date range",   BRAND),
    ("CSV Export",          "One-click download of raw entry data for external analysis", EMERALD),
    ("Access Control",      "Members see own data only · Admins & PMs see full team",     VIOLET),
    ("Visualizations",      "Daily bar chart + project time pie chart (Recharts)",        AMBER),
]
for i, (title, body, col) in enumerate(feats):
    c = i % 2
    r = i // 2
    x = 0.5 + c * 6.35
    y = 4.75 + r * 1.12
    txt(s, title, x,      y,        3.0, 0.38, size=12, bold=True, color=col)
    txt(s, body,  x,      y+0.42,   6.0, 0.55, size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLIENTS, SETTINGS & PROFILE
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
section_header(s, 7, "Clients · Settings · Profile")
txt(s, "Client CRM with logo storage · workspace configuration · personal avatar & password management",
    0.45, 1.25, 12.5, 0.38, size=12, color=MUTED)

# Clients
rect(s, 0.35, 1.6, 4.1, 5.7, SURFACE)
txt(s, "Clients", 0.6, 1.75, 3.6, 0.42, size=14, bold=True, color=WHITE)
for i, f in enumerate(["Name, email, notes", "Logo upload (PNG/JPG/SVG, 2 MB)", "Color avatar fallback",
                        "Full billing address", "VAT ID storage", "Project count display",
                        "Logo used on invoices", "Soft-delete (data preserved)"]):
    txt(s, f"✓  {f}", 0.6, 2.32+i*0.62, 3.7, 0.52, size=11, color=WHITE)

# Settings
rect(s, 4.75, 1.6, 4.1, 5.7, SURFACE)
txt(s, "Workspace Settings", 5.0, 1.75, 3.7, 0.42, size=14, bold=True, color=WHITE)
for i, f in enumerate(["Invite members by email", "Manage consultant levels", "Set weekly hours / member",
                        "Company legal name & address", "VAT ID, company reg number",
                        "IBAN / BIC for invoice terms", "Revenue & debitor accounts", "BMD NTCS tax code config"]):
    txt(s, f"✓  {f}", 5.0, 2.32+i*0.62, 3.7, 0.52, size=11, color=WHITE)

# Profile
rect(s, 9.15, 1.6, 3.85, 5.7, SURFACE)
txt(s, "Profile & Auth", 9.4, 1.75, 3.4, 0.42, size=14, bold=True, color=WHITE)
for i, f in enumerate(["Full name & email", "Avatar upload (Supabase storage)", "Live sidebar update on save",
                        "Password change + confirm", "Dark / Light / System theme",
                        "Language: DE / EN", "Invite-code gated signup", "Rate-limited auth (5/15 min)"]):
    txt(s, f"✓  {f}", 9.4, 2.32+i*0.62, 3.5, 0.52, size=11, color=WHITE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — TECHNICAL STACK
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Built to Last", 0.55, 0.28, 9, 0.6, size=30, bold=True, color=WHITE)
txt(s, "Production-grade stack · security-first · compliance-ready · mobile-responsive",
    0.55, 0.93, 12.5, 0.38, size=13, color=MUTED)

tech = [
    ("Next.js 15 + React 19",    "Latest App Router · server components · async cookies · streaming",         BRAND),
    ("Supabase / PostgreSQL",    "Row-Level Security on every table · service-role API for sensitive ops",     EMERALD),
    ("Auth & Rate Limiting",     "Invite-code gated · 5 attempts / 15 min IP limit · magic-link invites",     ROSE),
    ("Recharts Visualizations",  "Area, bar, and pie charts · fully responsive · dark-mode aware",             VIOLET),
    ("Dark / Light Mode",        "Full dark-mode support via next-themes · all components tested both ways",   MUTED),
    ("Bilingual EN / DE",        "Custom i18n system · locale-aware date & currency formatting throughout",    AMBER),
    ("Mobile PWA",               "Tailwind CSS · bottom nav on mobile · touch targets · service worker",       CYAN),
    ("Proxy / Impersonate",      "Admins view workspace exactly as any member — for support & QA",             PINK),
    ("Real-Time Presence",       "Supabase channels broadcast live 'who is tracking now' to all clients",      BRAND),
]
for i, (title, body, col) in enumerate(tech):
    c = i % 3
    r = i // 3
    x = 0.35 + c * 4.32
    y = 1.55 + r * 1.85
    card(s, x, y, 4.0, 1.68, col)
    txt(s, title, x+0.2, y+0.12, 3.6, 0.42, size=12, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.65, 3.6, 0.88, size=10, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — DIFFERENTIATORS
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Why Kairos?", 0.55, 0.28, 9, 0.6, size=30, bold=True, color=WHITE)
txt(s, "What sets Kairos apart from generic time-tracking tools",
    0.55, 0.93, 12.5, 0.38, size=13, color=MUTED)

diffs = [
    ("End-to-End Consulting Workflow",
     "Not a generic timer. Kairos covers the full lifecycle: time entry → PM review → "
     "partner sign-off → invoicing → BMD accounting export. Every step in one product.", BRAND),
    ("Austrian Compliance Out of the Box",
     "EN 16931 invoice structure, all Austrian VAT codes (U20/U10/IG/RC/AU), "
     "BMD NTCS export with configurable Erlöskonto & Debitorenkonto.", AMBER),
    ("Per-Project PM Approval Authority",
     "Each PM's approval rights are scoped to exactly their projects. "
     "The right person approves the right hours — enforced in UI and API.", EMERALD),
    ("Approved-Hours-Only Invoicing",
     "You can only invoice formally approved hours. "
     "No more billing unreviewed or disputed entries.", VIOLET),
    ("Burnout & Budget Alerts",
     "Flags >100 % utilization for 3+ consecutive weeks. "
     "Predicts project budget exhaustion based on current daily burn rate.", ROSE),
    ("Real-Time Team Visibility",
     "See who is tracking right now, for how long, on which project — "
     "without asking anyone. Live presence updates across all clients.", CYAN),
]
for i, (title, body, col) in enumerate(diffs):
    c = i % 2
    r = i // 2
    x = 0.35 + c * 6.55
    y = 1.55 + r * 1.9
    card(s, x, y, 6.1, 1.72, col)
    txt(s, title, x+0.2, y+0.12, 5.7, 0.42, size=13, bold=True, color=WHITE)
    txt(s, body,  x+0.2, y+0.65, 5.7, 0.9,  size=11, color=MUTED)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — IDEAL CUSTOMER
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
accent_bar(s)
txt(s, "Who Is Kairos For?", 0.55, 0.28, 9, 0.6, size=30, bold=True, color=WHITE)
txt(s, "Purpose-built for professional service firms that bill by the hour",
    0.55, 0.93, 12.5, 0.38, size=13, color=MUTED)

personas = [
    ("Management Consultancies",  "5–50 consultants · multiple clients · structured approval layers · Austrian VAT compliance",  BRAND),
    ("Law & Tax Advisory Firms",  "Billable hour tracking · partner approval · client-level invoicing · IBAN on invoice",        VIOLET),
    ("IT & Software Agencies",    "Project-based billing · level-based rates · PM-scoped review · budget burndown alerts",       EMERALD),
    ("Engineering Consultancies", "Multi-project teams · daily rate billing · timesheet locking · compliance-ready exports",     AMBER),
]
for i, (title, body, col) in enumerate(personas):
    card(s, 0.35, 1.55+i*1.38, 12.6, 1.22, col)
    txt(s, title, 0.6, 1.7+i*1.38,  4.5, 0.42, size=14, bold=True, color=WHITE)
    txt(s, body,  5.3, 1.72+i*1.38, 7.5, 0.42, size=11, color=MUTED)

rect(s, 0.35, 7.1, 12.6, 0.3, RGBColor(0x10, 0x0B, 0x3A))
txt(s, "Sweet spot:  5–50 person professional service firm  ·  hourly or daily billing  ·  Austrian / German market",
    0.5, 7.12, 12.2, 0.26, size=10, color=INDIGO, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — CLOSING / CTA
# ══════════════════════════════════════════════════════════════════════════════
s = add_slide()
bg(s)
rect(s, 0, 0, 5.5, 7.5, SURFACE)
rect(s, 0, 0, 0.08, 7.5, BRAND)

txt(s, "KAIROS", 0.5, 1.4, 4.8, 1.0, size=52, bold=True, color=WHITE)
txt(s, "Consulting Management Platform", 0.5, 2.5, 5.0, 0.42, size=13, color=MUTED)
rect(s, 0.5, 3.05, 4.7, 0.055, BRAND)
txt(s, "Track every hour.\nApprove with authority.\nInvoice with confidence.",
    0.5, 3.22, 5.0, 1.1, size=14, italic=True, color=INDIGO)
txt(s, "Built in Austria  ·  EN 16931  ·  BMD NTCS ready",
    0.5, 4.55, 5.0, 0.35, size=11, color=MUTED)
txt(s, "Contact for Demo & Pricing", 0.5, 5.45, 5.0, 0.38, size=13, bold=True, color=WHITE)
txt(s, "kairos.consulting", 0.5, 5.9, 5.0, 0.35, size=12, color=MUTED)

summary = [
    ("9 Modules",       "Timer to BMD export",            BRAND),
    ("4 Roles",         "Admin, Partner, PM, Member",      VIOLET),
    ("EN 16931",        "Invoice standard compliance",     EMERALD),
    ("BMD NTCS",        "Direct accounting export",        AMBER),
    ("Real-Time",       "Live team presence tracking",     ROSE),
    ("DE / EN",         "Fully bilingual interface",       CYAN),
    ("PWA Ready",       "Mobile-first responsive design",  VIOLET),
    ("RLS Security",    "Row-Level Security throughout",   EMERALD),
]
txt(s, "What you get", 5.8, 0.85, 7.2, 0.42, size=15, bold=True, color=WHITE)
for i, (val, label, col) in enumerate(summary):
    c = i % 2
    r = i // 2
    x = 5.8 + c * 3.7
    y = 1.45 + r * 1.42
    card(s, x, y, 3.4, 1.25, col)
    txt(s, val,   x+0.2, y+0.1,  3.0, 0.5,  size=16, bold=True, color=col)
    txt(s, label, x+0.2, y+0.68, 3.0, 0.42, size=10, color=MUTED)

out = "c:/Users/Max/OneDrive/Desktop/Kairos_Product_Deck.pptx"
prs.save(out)
print(f"Saved: {out}")
